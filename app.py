import json
import os
import random
import re
import tempfile
from dataclasses import dataclass
from datetime import datetime, timezone
from typing import Any, Dict, List, Optional, Tuple

import streamlit as st
from dotenv import load_dotenv
from openai import OpenAI
from pydantic import BaseModel, ValidationError
from pypdf import PdfReader
from docx import Document
from pptx import Presentation

load_dotenv()


# -----------------------------
# Data models
# -----------------------------
class QuizQuestion(BaseModel):
    """One multiple-choice question with answer key and concept tag."""
    question: str
    options: List[str]
    correct_index: int
    explanation: str
    concept: str


class StudyPack(BaseModel):
    """Bundle returned by generation pipeline: summary, concepts, and quiz."""
    summary: str
    key_concepts: List[str]
    quiz: List[QuizQuestion]


class ExplanationsPack(BaseModel):
    """Post-quiz remediation content for incorrect responses."""
    explanations: List[str]
    recommendations: List[str]


@dataclass
class QuizResult:
    """Normalized quiz evaluation output used by reporting and routing."""
    score: int
    total: int
    accuracy: float
    wrong_indices: List[int]
    confidence_mismatch: bool
    avg_confidence: float
    next_difficulty: str


# -----------------------------
# Helpers
# -----------------------------
API_KEY_FILE = os.getenv("OPENAI_API_KEY_FILE", "secrets/openai_api_key.txt")
MAX_CHARS_PER_CHUNK = 10000
MEMORY_FILE = "data/user_memory.json"


def normalize_topic(topic: str) -> str:
    """Create canonical topic key for stable per-topic memory lookup."""
    return re.sub(r"\s+", " ", topic.strip().lower())


def load_memory(path: str = MEMORY_FILE) -> Dict[str, Any]:
    """Load persistent learner memory from disk.

    The memory file stores topic-level aggregates and per-attempt quiz history.
    If the file does not exist or cannot be parsed, this returns a safe empty
    structure so the app can continue without failure.
    """
    if not os.path.exists(path):
        return {"topics": {}}
    try:
        with open(path, "r", encoding="utf-8") as f:
            data = json.load(f)
        if not isinstance(data, dict):
            return {"topics": {}}
        data.setdefault("topics", {})
        return data
    except (json.JSONDecodeError, OSError):
        return {"topics": {}}


def save_memory(memory: Dict[str, Any], path: str = MEMORY_FILE) -> None:
    """Persist learner memory to disk as formatted JSON."""
    os.makedirs(os.path.dirname(path), exist_ok=True)
    with open(path, "w", encoding="utf-8") as f:
        json.dump(memory, f, indent=2)


def clear_memory(path: str = MEMORY_FILE) -> None:
    """Delete persistent learner memory file if it exists."""
    if os.path.exists(path):
        os.remove(path)


def get_topic_record(memory: Dict[str, Any], topic: str) -> Dict[str, Any]:
    """Return mutable memory record for a topic, creating it on first use."""
    topic_key = normalize_topic(topic)
    topics = memory.setdefault("topics", {})
    if topic_key not in topics:
        topics[topic_key] = {
            "display_topic": topic.strip(),
            "sessions": 0,
            "last_avg_confidence": None,
            "last_routed_difficulty": None,
            "concept_stats": {},
            "quiz_history": [],
        }
    return topics[topic_key]


def confidence_to_difficulty(avg_confidence: float) -> str:
    """Map average confidence to next quiz difficulty.

    Rule:
    - 1-2 => foundational
    - 3   => standard
    - 4-5 => advanced
    """
    rounded = int(round(avg_confidence))
    if rounded <= 2:
        return "foundational"
    if rounded == 3:
        return "standard"
    return "advanced"


def top_mistake_concepts(topic_record: Dict[str, Any], limit: int = 5) -> List[str]:
    """Rank concepts by mistake frequency to prioritize future quiz generation."""
    concept_stats = topic_record.get("concept_stats", {})
    ranked = sorted(
        concept_stats.items(),
        key=lambda kv: (kv[1].get("wrong", 0), kv[1].get("seen", 0)),
        reverse=True,
    )
    return [concept for concept, stat in ranked if stat.get("wrong", 0) > 0][:limit]


def update_topic_memory(
    memory: Dict[str, Any],
    topic: str,
    quiz: List[QuizQuestion],
    wrong_indices: List[int],
    avg_confidence: float,
    routed_difficulty: str,
    score: int,
    total: int,
    accuracy: float,
    confidence_mismatch: bool,
) -> None:
    """Update per-topic statistics and append one per-attempt history record."""
    rec = get_topic_record(memory, topic)
    rec["sessions"] = int(rec.get("sessions", 0)) + 1
    rec["last_avg_confidence"] = avg_confidence
    rec["last_routed_difficulty"] = routed_difficulty
    concept_stats = rec.setdefault("concept_stats", {})

    wrong_set = set(wrong_indices)
    for i, q in enumerate(quiz):
        c = q.concept.strip().lower()
        stat = concept_stats.setdefault(c, {"seen": 0, "wrong": 0})
        stat["seen"] += 1
        if i in wrong_set:
            stat["wrong"] += 1

    wrong_concepts = [quiz[i].concept.strip().lower() for i in wrong_indices]
    history = rec.setdefault("quiz_history", [])
    history.append(
        {
            "timestamp_utc": datetime.now(timezone.utc).isoformat(),
            "score": score,
            "total": total,
            "accuracy_pct": round(accuracy * 100, 1),
            "avg_confidence": round(avg_confidence, 2),
            "routed_difficulty": routed_difficulty,
            "confidence_mismatch": confidence_mismatch,
            "wrong_concepts": wrong_concepts,
        }
    )
    # Keep latest 100 entries per topic to limit file growth.
    rec["quiz_history"] = history[-100:]


def read_api_key_from_file(path: str) -> Optional[str]:
    """Read API key from local text file; return None if missing/empty."""
    if not os.path.exists(path):
        return None
    with open(path, "r", encoding="utf-8") as f:
        key = f.read().strip()
    return key or None


def get_openai_client() -> Optional[OpenAI]:
    """Build OpenAI client from env key or local key file."""
    api_key = os.getenv("OPENAI_API_KEY") or read_api_key_from_file(API_KEY_FILE)
    if not api_key:
        return None
    return OpenAI(api_key=api_key)


def extract_text_from_file(uploaded_file) -> str:
    """Extract text for supported file types and clean up temp parser files."""
    file_name = uploaded_file.name.lower()

    with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(file_name)[1]) as tmp:
        tmp.write(uploaded_file.getvalue())
        temp_path = tmp.name

    try:
        if file_name.endswith(".txt"):
            return uploaded_file.getvalue().decode("utf-8", errors="ignore")

        if file_name.endswith(".pdf"):
            reader = PdfReader(temp_path)
            return "\n".join((page.extract_text() or "") for page in reader.pages)

        if file_name.endswith(".docx"):
            doc = Document(temp_path)
            return "\n".join(p.text for p in doc.paragraphs)

        if file_name.endswith(".pptx"):
            prs = Presentation(temp_path)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
            return "\n".join(texts)

        raise ValueError("Unsupported file type. Please upload TXT, PDF, DOCX, or PPTX.")
    finally:
        if os.path.exists(temp_path):
            os.remove(temp_path)


def clean_text(text: str) -> str:
    """Normalize whitespace to reduce prompt noise and parsing edge cases."""
    text = re.sub(r"\s+", " ", text).strip()
    return text


def chunk_text(text: str, max_chars: int = MAX_CHARS_PER_CHUNK) -> List[str]:
    """Split long text into fixed-size chunks for multi-call LLM processing."""
    text = clean_text(text)
    if len(text) <= max_chars:
        return [text]
    chunks = []
    start = 0
    while start < len(text):
        end = min(start + max_chars, len(text))
        chunks.append(text[start:end])
        start = end
    return chunks


def safe_json_load(raw: str) -> Dict[str, Any]:
    """Parse JSON responses and tolerate fenced markdown wrappers."""
    raw = raw.strip()
    if raw.startswith("```"):
        raw = re.sub(r"^```(?:json)?", "", raw).strip()
        raw = re.sub(r"```$", "", raw).strip()
    return json.loads(raw)


def sentence_chunks(text: str, n: int = 5) -> List[str]:
    """Return deterministic sentence slices for fallback summary/question seeds."""
    sentences = re.split(r"(?<=[.!?])\s+", text)
    sentences = [s.strip() for s in sentences if len(s.strip()) > 20]
    if not sentences:
        return ["No substantial sentence found in source material."] * n
    return [sentences[i % len(sentences)] for i in range(n)]


def weighted_sample_concepts(
    concepts: List[str], focus_concepts: Optional[List[str]], k: int = 5
) -> List[str]:
    """Sample quiz concepts with replacement weighted toward historically weak concepts."""
    if not concepts:
        return []

    unique = list(dict.fromkeys([c.strip().lower() for c in concepts if c.strip()]))
    if not unique:
        return []

    focus_set = set(c.strip().lower() for c in (focus_concepts or []) if c.strip())
    weights = [3.0 if c in focus_set else 1.0 for c in unique]

    # weighted random picks; de-duplicate while preserving draw order.
    drawn = random.choices(unique, weights=weights, k=min(k * 3, max(10, k)))
    picked = list(dict.fromkeys(drawn))

    # Ensure exactly k concepts when possible.
    for c in unique:
        if len(picked) >= min(k, len(unique)):
            break
        if c not in picked:
            picked.append(c)
    return picked[: min(k, len(unique))]


def fallback_study_pack(
    text: str, difficulty: str = "standard", focus_concepts: Optional[List[str]] = None
) -> StudyPack:
    """Offline deterministic generation path when API key/model call is unavailable."""
    chunks = sentence_chunks(text, 8)
    summary = " ".join(chunks[:3])

    concepts = []
    words = re.findall(r"[A-Za-z]{6,}", text)
    for w in words:
        lw = w.lower()
        if lw not in concepts:
            concepts.append(lw)
        if len(concepts) >= 40:
            break
    if not concepts:
        concepts = ["core concept", "key idea", "application", "analysis"]
    if focus_concepts:
        prioritized = [c.strip().lower() for c in focus_concepts if c.strip()]
        concepts = list(dict.fromkeys(prioritized + concepts))

    quiz_concepts = weighted_sample_concepts(concepts, focus_concepts, k=5)
    if len(quiz_concepts) < 5:
        quiz_concepts = (quiz_concepts + concepts)[:5]

    questions: List[QuizQuestion] = []
    level_word = {
        "foundational": "basic",
        "standard": "intermediate",
        "advanced": "advanced",
    }.get(difficulty, "intermediate")

    for i in range(5):
        concept = quiz_concepts[i % len(quiz_concepts)]
        stem = chunks[(i + 2) % len(chunks)]
        question = f"({level_word}) What is the best interpretation of this material segment? {stem[:120]}"
        if difficulty == "foundational":
            question += f" Hint: Focus on the core meaning of '{concept}'."
        options = [
            f"Option A: Definition related to {concept}",
            f"Option B: Misinterpretation of {concept}",
            f"Option C: Irrelevant detail",
            f"Option D: Opposite claim",
        ]
        questions.append(
            QuizQuestion(
                question=question,
                options=options,
                correct_index=0,
                explanation=f"The text emphasizes the core idea of {concept}, so the definition-focused option is best.",
                concept=concept,
            )
        )

    return StudyPack(summary=summary, key_concepts=concepts, quiz=questions)


def summarize_chunk_with_llm(client: OpenAI, chunk: str, model: str = "gpt-4o-mini") -> Dict[str, Any]:
    """Summarize one chunk and extract local concepts for map-reduce aggregation."""
    prompt = f"""
You are summarizing one chunk of study material.
Return strict JSON with keys:
- chunk_summary: string (2-4 sentences)
- key_concepts: array of up to 5 short concept strings
Rules:
- Use only the provided chunk content.
- Do not introduce facts, definitions, or concepts not present in the chunk.

Chunk:
{chunk}
"""
    response = client.responses.create(
        model=model,
        input=prompt,
        temperature=0.2,
    )
    return safe_json_load(response.output_text)


def generate_study_pack_with_llm(
    client: OpenAI,
    text: str,
    difficulty: str = "standard",
    model: str = "gpt-4o-mini",
    focus_concepts: Optional[List[str]] = None,
    topic: Optional[str] = None,
) -> StudyPack:
    """Create summary + concept inventory, then generate weighted-random concept quiz."""

    chunks = chunk_text(text)
    if len(chunks) == 1:
        source_for_generation = chunks[0]
    else:
        # Map-reduce style: summarize each chunk first, then generate one final
        # study pack from aggregated summaries and consolidated concept hints.
        chunk_summaries: List[str] = []
        concept_pool: List[str] = []
        for chunk in chunks:
            partial = summarize_chunk_with_llm(client, chunk, model)
            summary = clean_text(str(partial.get("chunk_summary", "")))
            if summary:
                chunk_summaries.append(summary)

            for concept in partial.get("key_concepts", []):
                c = clean_text(str(concept)).lower()
                if c and c not in concept_pool:
                    concept_pool.append(c)

        source_for_generation = (
            "Aggregated chunk summaries:\n"
            + "\n".join(f"- {s}" for s in chunk_summaries[:80])
            + "\n\nCandidate concepts:\n"
            + ", ".join(concept_pool[:30])
        )

    hint_instruction = (
        "For foundational difficulty, include a short hint in each question stem.\n"
        if difficulty == "foundational"
        else ""
    )

    topic_instruction = f"Topic label: {topic}\n" if topic else ""

    concept_prompt = f"""
Source material:
{source_for_generation}
{topic_instruction}

Task:
1) Create a concise structured summary suitable for student revision.
2) Extract a comprehensive list of key concepts from the source material (not limited to 5).
3) Use only the provided source material. Do not use external knowledge.
Return only valid JSON.
"""

    concept_response = client.responses.create(
        model=model,
        input=[
            {
                "role": "system",
                "content": (
                    "Return strict JSON with keys: summary (string), key_concepts (array of strings). "
                    "Use only provided source material. Do not invent content."
                ),
            },
            {"role": "user", "content": concept_prompt},
        ],
        temperature=0.2,
    )
    concept_data = safe_json_load(concept_response.output_text)
    summary = clean_text(str(concept_data.get("summary", "")))
    key_concepts_raw = concept_data.get("key_concepts", [])
    key_concepts = list(
        dict.fromkeys(
            clean_text(str(c)).lower()
            for c in key_concepts_raw
            if clean_text(str(c))
        )
    )
    if not key_concepts:
        key_concepts = ["core concept", "key idea", "application", "analysis"]

    selected_quiz_concepts = weighted_sample_concepts(key_concepts, focus_concepts, k=5)
    if len(selected_quiz_concepts) < 5:
        selected_quiz_concepts = (selected_quiz_concepts + key_concepts)[:5]

    quiz_prompt = f"""
Source material:
{source_for_generation}
{topic_instruction}

Target quiz concepts (generate one question per concept, in this order):
{json.dumps(selected_quiz_concepts)}

Task:
1) Create exactly 5 multiple-choice questions at {difficulty} difficulty.
2) Each question must be grounded in its assigned target concept from the provided list.
3) Questions must test understanding, not just copying text.
4) Keep explanations short but instructive.
5) For each question, explanation must include why correct option is correct and why a likely wrong option is wrong.
6) Every question and option must be fully answerable from the provided source material only.
7) Do not use external facts, assumptions, or prior domain knowledge.
{hint_instruction}
Return only valid JSON.
"""

    quiz_response = client.responses.create(
        model=model,
        input=[
            {
                "role": "system",
                "content": (
                    "Return strict JSON with key quiz (array of exactly 5 objects). "
                    "Each object: question (string), options (array of 4 strings), "
                    "correct_index (0-3 int), explanation (string), concept (string). "
                    "Use only provided source material."
                ),
            },
            {"role": "user", "content": quiz_prompt},
        ],
        temperature=0.3,
    )

    quiz_data = safe_json_load(quiz_response.output_text)
    quiz_items = quiz_data.get("quiz", [])
    return StudyPack.model_validate(
        {"summary": summary, "key_concepts": key_concepts, "quiz": quiz_items}
    )


def evaluate_quiz(
    quiz: List[QuizQuestion],
    answers: List[int],
    confidence: List[int],
) -> QuizResult:
    """Score quiz responses and compute confidence-driven next difficulty route."""
    score = 0
    wrong_indices: List[int] = []
    mismatch_count = 0

    for i, q in enumerate(quiz):
        chosen = answers[i]
        is_correct = chosen == q.correct_index
        if is_correct:
            score += 1

        if not is_correct:
            wrong_indices.append(i)

        # Mismatch: high confidence but wrong, or low confidence but correct.
        if (confidence[i] >= 4 and not is_correct) or (confidence[i] <= 2 and is_correct):
            mismatch_count += 1

    avg_conf = sum(confidence) / len(confidence)
    next_diff = confidence_to_difficulty(avg_conf)

    return QuizResult(
        score=score,
        total=len(quiz),
        accuracy=score / len(quiz),
        wrong_indices=wrong_indices,
        confidence_mismatch=mismatch_count >= 2,
        avg_confidence=avg_conf,
        next_difficulty=next_diff,
    )


def generate_explanations_with_llm(
    client: OpenAI,
    quiz: List[QuizQuestion],
    answers: List[int],
    wrong_indices: List[int],
    key_concepts: List[str],
    model: str = "gpt-4o-mini",
) -> ExplanationsPack:
    """Generate personalized explanations for incorrect responses."""
    mistakes = []
    for idx in wrong_indices:
        q = quiz[idx]
        mistakes.append(
            {
                "question": q.question,
                "chosen_option": q.options[answers[idx]],
                "correct_option": q.options[q.correct_index],
                "concept": q.concept,
            }
        )

    prompt = f"""
The learner answered these questions incorrectly:
{json.dumps(mistakes, indent=2)}

Key concepts list:
{json.dumps(key_concepts)}

Provide strict JSON with:
- explanations: array of short personalized explanations aligned to each mistake.
  Each explanation must explicitly include:
  (a) why the correct option is correct
  (b) why the learner's selected option is incorrect
  Format each explanation with readable markdown sections:
  **Concept:** ...
  **Your Answer:** ...
  **Correct Answer:** ...
  **Why Correct:** ...
  **Why Your Answer Is Incorrect:** ...
- recommendations: array of 3 targeted study actions
"""

    response = client.responses.create(
        model=model,
        input=prompt,
        temperature=0.2,
    )
    parsed = safe_json_load(response.output_text)
    return ExplanationsPack.model_validate(parsed)


def fallback_explanations(
    quiz: List[QuizQuestion], answers: List[int], wrong_indices: List[int], key_concepts: List[str]
) -> ExplanationsPack:
    """Deterministic explanation fallback when LLM output is unavailable."""
    exps = []
    for idx in wrong_indices:
        q = quiz[idx]
        chosen = q.options[answers[idx]]
        correct = q.options[q.correct_index]
        exps.append(
            f"**Concept:** {q.concept}\n\n"
            f"**Your Answer:** {chosen}\n\n"
            f"**Correct Answer:** {correct}\n\n"
            f"**Why Correct:** It matches the key idea from the uploaded material.\n\n"
            f"**Why Your Answer Is Incorrect:** It does not align with the concept focus in the material.\n\n"
            f"**Extra Detail:** {q.explanation}"
        )

    recs = [
        f"Review concept map for: {', '.join(key_concepts[:3])}.",
        "Redo one short quiz immediately after revising mistakes.",
        "Write a 3-sentence explanation for each missed concept from memory.",
    ]
    return ExplanationsPack(explanations=exps, recommendations=recs)


def init_state() -> None:
    """Initialize all Streamlit session keys required by the app."""
    # Session state stores transient UI/session data; persistent learner history is
    # saved separately in data/user_memory.json.
    defaults = {
        "topic": "",
        "source_text": "",
        "source_file_count": 0,
        "study_pack": None,
        "next_quiz_pack": None,
        "quiz_submitted": False,
        "result": None,
        "explanations": None,
        "active_focus_concepts": [],
        "active_difficulty": "standard",
        "quiz_attempt_number": 1,
        "uploader_key_version": 0,
    }
    for k, v in defaults.items():
        if k not in st.session_state:
            st.session_state[k] = v


def clear_current_outputs(preserve_topic: bool = True) -> None:
    """Reset generated UI outputs while keeping optional topic text in the input box."""
    # Reset all generated content and quiz state while optionally preserving topic input.
    topic_value = st.session_state.topic if preserve_topic else ""
    st.session_state.source_text = ""
    st.session_state.source_file_count = 0
    st.session_state.study_pack = None
    st.session_state.next_quiz_pack = None
    st.session_state.quiz_submitted = False
    st.session_state.result = None
    st.session_state.explanations = None
    st.session_state.active_focus_concepts = []
    st.session_state.active_difficulty = "standard"
    st.session_state.quiz_attempt_number = 1
    st.session_state.uploader_key_version += 1
    st.session_state.topic = topic_value


# -----------------------------
# UI
# -----------------------------
# Main page layout: setup controls, topic/material intake, generation, quiz interaction,
# and side-by-side reporting once quiz results are available.
st.set_page_config(page_title="Adaptive Study Agent", layout="wide")
st.title("Academic Tutor & Adaptive Quiz Agent")
st.caption(
    "Uploads study material, generates a concept-focused quiz, evaluates responses, "
    "and adapts future quizzes based on confidence and previous mistakes."
)

init_state()
client = get_openai_client()
llm_mode = "OpenAI API" if client else "Fallback (No API key detected)"
memory = load_memory()

with st.sidebar:
    st.subheader("Settings")
    selected_model = st.text_input("Model", value="gpt-4o-mini")
    st.write(f"Mode: **{llm_mode}**")
    st.info(
        f"Use `OPENAI_API_KEY` env var or store key in `{API_KEY_FILE}` to use LLM mode."
    )
    st.caption(f"Memory file: `{MEMORY_FILE}`")
    with st.popover("Clear Current Page Output"):
        st.write(
            "This clears generated summaries, quizzes, reports, and uploaded files from "
            "the current page. The topic text will be kept."
        )
        if st.button("Confirm Clear Current Output", type="secondary"):
            clear_current_outputs(preserve_topic=True)
            st.rerun()
    with st.popover("Clear All Saved Memory"):
        st.write(
            "This permanently removes all saved topic memory from local storage "
            f"(`{MEMORY_FILE}`). Current page output and uploaded files will also be cleared."
        )
        if st.button("Confirm Clear All Saved Memory", type="secondary"):
            clear_memory()
            clear_current_outputs(preserve_topic=True)
            st.rerun()

st.session_state.topic = st.text_input(
    "Topic for this session (used for memory and adaptive future quizzes)",
    value=st.session_state.topic,
    placeholder="e.g., Cellular Respiration",
)
active_topic = st.session_state.topic.strip()
if active_topic:
    existing = memory.get("topics", {}).get(normalize_topic(active_topic), {})
    if existing:
        st.caption(
            f"Past sessions for this topic: {existing.get('sessions', 0)} | "
            f"Last avg confidence: {existing.get('last_avg_confidence', 'N/A')}"
        )

uploaded = st.file_uploader(
    "Upload study material (TXT, PDF, DOCX, PPTX)",
    type=["txt", "pdf", "docx", "pptx"],
    accept_multiple_files=True,
    key=f"uploader_{st.session_state.uploader_key_version}",
)

if uploaded:
    # Combine extracted text from all selected files into one source corpus so
    # downstream summarization and quiz generation can operate across documents.
    all_texts: List[str] = []
    failed_files: List[Tuple[str, str]] = []
    for f in uploaded:
        try:
            extracted = clean_text(extract_text_from_file(f))
            if extracted:
                labeled = f"[FILE: {f.name}]\n{extracted}"
                all_texts.append(labeled)
        except Exception as e:
            failed_files.append((f.name, str(e)))

    full_text = "\n\n".join(all_texts)
    st.session_state.source_text = full_text
    st.session_state.source_file_count = len(all_texts)

    if all_texts:
        st.success(
            f"Loaded {len(all_texts)} file(s), total extracted length: {len(full_text):,} chars"
        )
        with st.expander("Preview extracted text", expanded=False):
            st.write(full_text[:3000] + ("..." if len(full_text) > 3000 else ""))

    for file_name, err in failed_files:
        st.error(f"Could not process file '{file_name}': {err}")

if st.button("Generate Summary + Initial Quiz", type="primary"):
    if not active_topic:
        st.warning("Please enter a topic first.")
    elif not st.session_state.source_text:
        st.warning("Please upload a document first.")
    else:
        # Use existing topic memory to bias the new quiz toward previously weak
        # concepts and to choose a starting difficulty for this session.
        topic_rec = get_topic_record(memory, active_topic)
        historical_focus = top_mistake_concepts(topic_rec, limit=5)
        prior_conf = topic_rec.get("last_avg_confidence")
        difficulty = (
            confidence_to_difficulty(prior_conf)
            if isinstance(prior_conf, (int, float))
            else "standard"
        )
        st.session_state.active_focus_concepts = historical_focus
        st.session_state.active_difficulty = difficulty

        if historical_focus:
            st.info(
                "Applying memory-weighted focus concepts for this topic: "
                + ", ".join(historical_focus)
            )
        if isinstance(prior_conf, (int, float)):
            st.info(
                f"Using previous topic confidence ({prior_conf:.2f}) -> "
                f"start difficulty: {difficulty}"
            )

        chunk_count = len(chunk_text(st.session_state.source_text))
        if chunk_count > 1:
            st.info(
                f"Large input detected from {st.session_state.source_file_count} file(s). "
                f"Processing in {chunk_count} chunk(s)."
            )
        with st.spinner("Generating study pack..."):
            try:
                # Primary path: LLM generation with concept-aware weighted quiz concept sampling.
                if client:
                    pack = generate_study_pack_with_llm(
                        client,
                        st.session_state.source_text,
                        difficulty,
                        selected_model,
                        focus_concepts=historical_focus,
                        topic=active_topic,
                    )
                else:
                    pack = fallback_study_pack(
                        st.session_state.source_text,
                        difficulty,
                        focus_concepts=historical_focus,
                    )
                st.session_state.study_pack = pack
                st.session_state.quiz_submitted = False
                st.session_state.result = None
                st.session_state.explanations = None
                st.session_state.next_quiz_pack = None
                st.session_state.quiz_attempt_number = 1
            except (ValidationError, json.JSONDecodeError, Exception):
                # Reliability fallback: keep the UX functional even if model JSON is malformed.
                pack = fallback_study_pack(
                    st.session_state.source_text,
                    difficulty,
                    focus_concepts=historical_focus,
                )
                st.session_state.study_pack = pack
                st.session_state.quiz_submitted = False
                st.session_state.result = None
                st.session_state.explanations = None
                st.session_state.next_quiz_pack = None
                st.session_state.quiz_attempt_number = 1
                st.warning("Switched to fallback generation due to model output format issues.")

pack: Optional[StudyPack] = st.session_state.study_pack
if pack:
    # Three-column layout: summary (left), quiz (center), report/explanations (right).
    left_col, center_col, right_col = st.columns([1.05, 1.35, 1.1])

    with left_col:
        st.subheader("Generated Summary")
        st.write(pack.summary)
        st.subheader("Key Concepts")
        st.write(", ".join(pack.key_concepts))

    with center_col:
        st.subheader(f"Current Quiz (Attempt {st.session_state.quiz_attempt_number})")
        with st.form(f"quiz_form_attempt_{st.session_state.quiz_attempt_number}"):
            answers: List[int] = []
            confidence: List[int] = []
            for i, q in enumerate(pack.quiz):
                st.markdown(f"**Q{i+1}. {q.question}**")
                choice = st.radio(
                    "Select one:",
                    options=list(range(4)),
                    format_func=lambda x, q=q: q.options[x],
                    key=f"q_{st.session_state.quiz_attempt_number}_{i}",
                )
                conf = st.slider(
                    "Confidence (1=guess, 5=very sure)",
                    min_value=1,
                    max_value=5,
                    value=3,
                    key=f"c_{st.session_state.quiz_attempt_number}_{i}",
                )
                answers.append(choice)
                confidence.append(conf)

            submitted = st.form_submit_button("Submit Quiz")

        if submitted:
            # Evaluate and persist the quiz attempt immediately so downstream
            # adaptation and history analytics stay in sync with user actions.
            result = evaluate_quiz(pack.quiz, answers, confidence)
            st.session_state.quiz_submitted = True
            st.session_state.result = result
            wrong_concepts_current = [pack.quiz[i].concept.strip().lower() for i in result.wrong_indices]

            if active_topic:
                # Persist this attempt for future topic-specific adaptation and
                # for the dedicated history page.
                update_topic_memory(
                    memory,
                    active_topic,
                    pack.quiz,
                    result.wrong_indices,
                    result.avg_confidence,
                    result.next_difficulty,
                    result.score,
                    result.total,
                    result.accuracy,
                    result.confidence_mismatch,
                )
                save_memory(memory)

            if result.wrong_indices:
                with st.spinner("Generating explanations..."):
                    try:
                        if client:
                            ex_pack = generate_explanations_with_llm(
                                client,
                                pack.quiz,
                                answers,
                                result.wrong_indices,
                                pack.key_concepts,
                                selected_model,
                            )
                        else:
                            ex_pack = fallback_explanations(
                                pack.quiz, answers, result.wrong_indices, pack.key_concepts
                            )
                    except Exception:
                        ex_pack = fallback_explanations(
                            pack.quiz, answers, result.wrong_indices, pack.key_concepts
                        )
                st.session_state.explanations = ex_pack

            # Generate adaptive second-round quiz as optional next step.
            with st.spinner("Preparing adaptive next round..."):
                try:
                    # Prioritize concepts that are weak in this attempt and in
                    # historical topic performance.
                    topic_rec_after = get_topic_record(memory, active_topic) if active_topic else {}
                    historical_focus = top_mistake_concepts(topic_rec_after, limit=5)
                    weighted_focus = list(
                        dict.fromkeys(wrong_concepts_current + historical_focus + st.session_state.active_focus_concepts)
                    )[:5]
                    st.session_state.active_focus_concepts = weighted_focus

                    if client:
                        round2 = generate_study_pack_with_llm(
                            client,
                            st.session_state.source_text,
                            result.next_difficulty,
                            selected_model,
                            focus_concepts=weighted_focus,
                            topic=active_topic,
                        )
                    else:
                        round2 = fallback_study_pack(
                            st.session_state.source_text,
                            result.next_difficulty,
                            focus_concepts=weighted_focus,
                        )
                    st.session_state.next_quiz_pack = round2
                except Exception:
                    st.session_state.next_quiz_pack = fallback_study_pack(
                        st.session_state.source_text,
                        result.next_difficulty,
                        focus_concepts=st.session_state.active_focus_concepts,
                    )
    with right_col:
        if st.session_state.quiz_submitted and st.session_state.result:
            result: QuizResult = st.session_state.result
            st.subheader("Performance Report")

            st.metric("Score", f"{result.score}/{result.total}")
            st.metric("Accuracy", f"{result.accuracy * 100:.0f}%")
            st.metric("Average Confidence", f"{result.avg_confidence:.2f}/5")
            st.write(
                f"**Adaptive route (confidence-based):** next quiz difficulty is "
                f"**{result.next_difficulty}**"
            )
            st.caption("Routing rule: confidence 1-2 -> foundational, 3 -> standard, 4-5 -> advanced")

            if st.session_state.active_focus_concepts:
                st.write(
                    "**Concept weighting for future quizzes:** "
                    + ", ".join(st.session_state.active_focus_concepts)
                )

            if result.confidence_mismatch:
                st.warning(
                    "Confidence mismatch detected (confidence did not align with correctness). "
                    "Added deeper explanations and recommend reviewing fundamentals of missed concepts."
                )

            ex_pack: Optional[ExplanationsPack] = st.session_state.explanations
            if ex_pack and ex_pack.explanations:
                st.subheader("Targeted Explanations")
                for i, exp in enumerate(ex_pack.explanations, start=1):
                    st.markdown(f"**Explanation {i}**")
                    st.markdown(exp)
                    if i < len(ex_pack.explanations):
                        st.markdown("---")

                st.subheader("Recommended Study Actions")
                for i, rec in enumerate(ex_pack.recommendations, start=1):
                    st.write(f"{i}. {rec}")

            next_quiz: Optional[StudyPack] = st.session_state.next_quiz_pack
            if next_quiz:
                st.subheader(f"Next Quiz Preview (Attempt {st.session_state.quiz_attempt_number + 1})")
                st.caption("Generated using confidence-based routing and historical weak-concept weighting.")
                for i, q in enumerate(next_quiz.quiz, start=1):
                    st.write(f"Q{i}: {q.question}")
                    for opt in q.options:
                        st.write(f"- {opt}")
                if st.button("Start Next Quiz"):
                    st.session_state.study_pack = next_quiz
                    st.session_state.quiz_submitted = False
                    st.session_state.result = None
                    st.session_state.explanations = None
                    st.session_state.next_quiz_pack = None
                    st.session_state.quiz_attempt_number += 1
                    st.rerun()
        else:
            st.subheader("Performance Report")
            st.caption("Submit the current quiz to see evaluation and targeted explanations here.")
